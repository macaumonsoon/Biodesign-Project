#!/usr/bin/env python3
"""
Single canonical BDC deck + PDF for GitHub in-browser preview.
Builds from the BDC_Extinction_Archive_EN.pptx template so the final deck
inherits its dimensions / theme more closely than a blank presentation.
"""

from pathlib import Path

from fpdf import FPDF
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt

ROOT = Path(__file__).resolve().parent.parent
EXPORT = ROOT / "slides" / "export"
OUT_PPTX = EXPORT / "Archive_of_Extinction_Final_BDC2026.pptx"
OUT_PDF = EXPORT / "Archive_of_Extinction_Final_BDC2026.pdf"
TEMPLATE = EXPORT / "BDC_Extinction_Archive_EN.pptx"

MAIN_TITLE = "Archive of Extinction: AI Memorial for Lost Species"
SUBTITLE_LINES = [
    "Umwelt Hypothesis Dossiers / Sensory Time Capsule",
    "Biodesign Challenge 2026 | Convergent Life | Biodigital Excellence",
    "Macau University",
]

INSTITUTION = "Macau University"
MENTOR = "Atticus SIMS"
TEAM = [
    ("GUO XIAO YUE", "MC569254"),
    ("LIU JIA QUN", "MC569293"),
]


def ascii_pdf(text: str) -> str:
    """FPDF core fonts: Windows-1252-ish; keep ASCII for reliability."""
    return (
        text.replace("\u2014", "-")
        .replace("\u2013", "-")
        .replace("\u201c", '"')
        .replace("\u201d", '"')
        .replace("\u2019", "'")
        .replace("\u2192", "->")
        .replace("\u00b7", "-")
        .replace("\u2265", ">=")
        .replace("\u00d7", "x")
    )


def slide_plan():
    """Ordered content merged from the highlighted decks + planning docs."""
    s = []
    s.append(("title", MAIN_TITLE, "\n".join(SUBTITLE_LINES)))
    s.append(
        (
            "bullets",
            "Team & institution",
            [
                f"Institution: {INSTITUTION}",
                f"Mentor: {MENTOR}",
                f"{TEAM[0][0]} - {TEAM[0][1]}",
                f"{TEAM[1][0]} - {TEAM[1][1]}",
                "Program: Digital Art & AI Technology (MDes) | 2-person team",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "Why this project now",
            [
                "Extinction removes more than bodies: it erases routes, light-dark entrainment, social synchrony, and sensory fit to landscape.",
                "Public memory normalizes ecological absence into icons; science stays in papers while loss becomes abstract.",
                "Design question: how do we make rigorous biology feelable without pretending we fully know another species' Umwelt?",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "Project thesis",
            [
                "Archive of Extinction is a browser-first biodesign memorial that reconstructs a bounded sensory-temporal hypothesis space for extinct species.",
                "It is grounded in citable research, expressed through WebXR, controlled generative passes, and optional Web Audio.",
                "The experience closes with evidence literacy and branching ethics so digital recall never erases conservation urgency.",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "Scope & title logic",
            [
                "Main title foregrounds extinction and remembrance: this is not de-extinction marketing.",
                "“Umwelt Hypothesis Dossiers” signals that each species is a folder of testable claims, not a fantasy safari.",
                "“Sensory Time Capsule” emphasizes temporal niche - how a species lived in day, season, synchrony, and atmosphere.",
                "Locked scope: two deep dossiers only - woolly mammoth + thylacine.",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "Core architecture",
            [
                "Dossier: evidence cards, temporal niche, sensory proxies, extinction mechanism, constrained 360 scene.",
                "Mixer: seasonal layers, diel rhythm, coexistence density, harmony -> conflict as sonic / temporal retuning.",
                "Ethics console: evidence gate + branching trade-offs around conservation, law, ecosystem risk, and sovereignty.",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "Biological layer",
            [
                "Only traceable evidence enters the archive: tusk isotopes, comparative genomics, morphology as sensory proxy, historical ecology, and acoustic relatives.",
                "Every strong claim is tagged Cited, Interpolated / Modeled, or Speculative.",
                "Mammoth copy is locked to “circadian-related genes” until full-text PER2 confirmation exists.",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "Digital layer & design techniques",
            [
                "WebXR stack: A-Frame / Three.js with desktop fallback; minimal QR-linked physical anchor for summit reliability.",
                "Generative image strategy: control-first passes, procedural atmosphere, and palette / sun rules bounded by paleo proxies.",
                "Time scrubber, sonification / Web Audio, and on-screen citation tags turn atmosphere into argument rather than decoration.",
                "Reflection UI uses auditable heuristics, not opaque LLM outputs, to compare viewer input with ethics themes.",
            ],
        )
    )
    s.append(
        (
            "two_col",
            "Two species - two dossiers",
            {
                "left_h": "Woolly mammoth",
                "left": [
                    "Hero species for Arctic rhythm, mammoth-steppe atmosphere, and de-extinction pressure.",
                    "Lynch 2015 circadian-related genes; Lu 2010 reindeer analogy [Interpolated]; Zimov 2012 habitat pulse.",
                    "Additional anchors: Wooller 2021 tusk isotope mobility; Nogues-Bravo 2008 climate x human extinction context.",
                ],
                "right_h": "Thylacine",
                "right": [
                    "Structural echo for crepuscular vision, colonial memory, and cultural sovereignty.",
                    "Pozniak 2018 orbit / diel activity; Mass & Supin 2020 RGC methods [Interpolated]; Paddle 2000; Sleightholme & Campbell 2016.",
                    "T5 context: Clements 2025 plus Palawa-led / centered Rimmer 2020, Lehman 2013, and Schlunke 2025.",
                ],
            },
        )
    )
    s.append(
        (
            "bullets",
            "User journey & scene flow",
            [
                "Mammoth path: Scene_Bio_01 -> Scene_Environment_02 -> Scene_Landscape_01 -> Scene_Sensory_01.",
                "Thylacine path: Scene_Bio_02 -> Scene_POV_01 -> Scene_Context_01.",
                "Shared modules: UI_Indigenous_Context -> UI_Ethics_Fork -> UI_Reflection_Log.",
                "All three load-bearing modules ship in the MVP: Dossier -> Mixer -> Ethics.",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "Methods & evidence protocol",
            [
                "Claims are imported to structured cards / JSON; each block is labeled Cited / Modeled / Speculative.",
                "Science visuals can include ggplot2 stills for tusk isotope summaries and time-series evidence panels.",
                "Generative runs log model, date, and seed where relevant; DSP rules are documented rather than hidden.",
                "Process includes at least two evidence audits during production to remove unlabeled claims.",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "Ethics, colonial violence & ecological memory",
            [
                "Sherkow & Greely (2013): trade-offs between laboratory de-extinction, in situ conservation, legal novelty, and ecological liability.",
                "Thylacine loss is entangled with violence against Palawa peoples and Country in lutruwita / Tasmania.",
                "Interface protocol: Palawa-first voices where possible; author positionality is explicit; no tourism copy as TEK.",
                "The archive is remembrance and accountability, not resurrection spectacle.",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "Ideation & design process",
            [
                "1 Map skills -> 2 Intersection statements -> 3 Rubric stress-test -> 4 Biology sprint.",
                "5 Bio / digital split -> 6 Narrative draft -> 7 Prototype + deliverables -> 8 Build / polish / Q&A.",
                "This process is aligned with the course workbook, biodesign_cursor_agent.md, and BDC judging logic.",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "Production plan (4-week spine)",
            [
                "Week 1: literature grid, storyboard, WebXR shell, evidence audit #1.",
                "Week 2: dossier UI, time scrubber, environment logic, ethics drafts, methods / limitations slide.",
                "Week 3: integrate thylacine context, reflection panel, sonification, evidence audit #2.",
                "Week 4: polish, accessibility, P0 bugs only, 10-min talk, 5-min Q&A, trailer export.",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "Biodesign Challenge fit",
            [
                "Narrative: clear emotional and spatial journey supported by live digital demo.",
                "Concept: temporal niche + Umwelt + evidence literacy creates a distinct biodigital proposition.",
                "Context: biodiversity loss, de-extinction ethics, misinformation risk, and Indigenous sovereignty.",
                "Reflection: documented limitations, epistemic layers, prompt / model logging, user reflection.",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "Deliverables & summit package",
            [
                "10-minute oral presentation + 5-minute Q&A.",
                "Visual renderings, WebXR experience, and citation-aware UI captures.",
                "Minimal physical anchor + QR for reliable access in gallery or talk format.",
                "1-5 minute creative trailer, PPT deck, methods appendix, and citations / limits in parallel.",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "Key literature - mammoth dossier",
            [
                "Lynch et al. 2015. Cell Reports. Elephantid genomes and Arctic adaptation. DOI 10.1016/j.celrep.2015.06.027",
                "Lu et al. 2010. Current Biology. Reindeer circadian attenuation. DOI 10.1016/j.cub.2010.01.042",
                "Zimov et al. 2012. Quaternary Science Reviews. Mammoth steppe productivity. DOI 10.1016/j.quascirev.2012.08.004",
                "Wooller et al. 2021. Science. Lifetime mobility of an Arctic woolly mammoth. DOI 10.1126/science.abg1134",
                "Nogues-Bravo et al. 2008. PLoS Biology. Climate, humans, and mammoth extinction. DOI 10.1371/journal.pbio.0060079",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "Key literature - thylacine dossier & sovereignty",
            [
                "Pozniak et al. 2018. Orbit / skull morphology and diel activity. DOI 10.1002/ar.23910",
                "Mass & Supin 2020. RGC topography methods. DOI 10.31857/S0002332920060107 [Interpolated in thylacine POV]",
                "Paddle 2000. The Last Tasmanian Tiger. ISBN 9780521782196",
                "Sleightholme & Campbell 2016. Human pressures and extinction framing. DOI 10.1080/00222933.2016.1217037",
                "Clements 2025. DOI 10.1177/13534858251272203 | Rimmer 2020 Artlink | Lehman 2013 Griffith Review | Schlunke 2025 DOI 10.1017/ext.2025.10008",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "Limitations & honest reflection",
            [
                "No living mammoth or thylacine: all sensory reconstruction is partial, bounded, and explicitly tiered.",
                "Generative media can feel real, so the deck and UI design against deceptive realism.",
                "Two people / short sprint means narrow but deep: one strong integrated experience beats many weak scenes.",
                "Desktop WebXR path is canonical; any AR or mobile path is supportive rather than mission-critical.",
            ],
        )
    )
    s.append(
        (
            "bullets",
            "Q&A prompts we practice",
            [
                "Why these species and not easier icons?",
                "Where is the living biology if there is no wet-lab organism on stage?",
                "How do you prevent AI misinformation and overclaiming?",
                "What would the next four months add: more dossiers, deeper methods, richer physical anchor, or consultant validation?",
            ],
        )
    )
    s.append(("title", "Thank you", "Archive of Extinction\nQuestions?"))
    return s


def remove_all_slides(prs):
    while len(prs.slides):
        slide_id = prs.slides._sldIdLst[0]
        rel_id = slide_id.get(qn("r:id"))
        prs.part.drop_rel(rel_id)
        prs.slides._sldIdLst.remove(slide_id)


def build_pptx(slides):
    prs = Presentation(str(TEMPLATE))
    remove_all_slides(prs)

    for item in slides:
        kind = item[0]
        if kind == "title":
            _, title, sub = item
            sl = prs.slides.add_slide(prs.slide_layouts[0])
            sl.shapes.title.text = title
            sl.placeholders[1].text = sub
        elif kind == "bullets":
            _, title, bullets = item
            sl = prs.slides.add_slide(prs.slide_layouts[1])
            sl.shapes.title.text = title
            tf = sl.placeholders[1].text_frame
            tf.clear()
            for i, b in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = b
                p.level = 0
                p.font.size = Pt(16)
        elif kind == "two_col":
            _, title, d = item
            sl = prs.slides.add_slide(prs.slide_layouts[3])
            sl.shapes.title.text = title
            bodies = [p for p in sl.placeholders if p.placeholder_format.idx > 0]
            if len(bodies) < 2:
                flat = [d["left_h"]] + d["left"] + [""] + [d["right_h"]] + d["right"]
                sl = prs.slides.add_slide(prs.slide_layouts[1])
                sl.shapes.title.text = title
                tf = sl.placeholders[1].text_frame
                tf.clear()
                for i, line in enumerate(flat):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = line
                    p.font.size = Pt(14)
                continue
            tf_l = bodies[0].text_frame
            tf_l.clear()
            tf_l.text = d["left_h"]
            tf_l.paragraphs[0].font.bold = True
            tf_l.paragraphs[0].font.size = Pt(14)
            for b in d["left"]:
                p = tf_l.add_paragraph()
                p.text = b
                p.font.size = Pt(13)
            tf_r = bodies[1].text_frame
            tf_r.clear()
            tf_r.text = d["right_h"]
            tf_r.paragraphs[0].font.bold = True
            tf_r.paragraphs[0].font.size = Pt(14)
            for b in d["right"]:
                p = tf_r.add_paragraph()
                p.text = b
                p.font.size = Pt(13)

    EXPORT.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))


def build_pdf(slides):
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_margins(left=18, top=18, right=18)

    def w():
        return pdf.epw

    for item in slides:
        kind = item[0]
        pdf.add_page()
        if kind == "title":
            _, title, sub = item
            pdf.set_font("Helvetica", "B", 22)
            pdf.multi_cell(w(), 10, ascii_pdf(title))
            pdf.ln(4)
            pdf.set_font("Helvetica", "", 12)
            pdf.multi_cell(w(), 6, ascii_pdf(sub))
        elif kind == "bullets":
            _, title, bullets = item
            pdf.set_font("Helvetica", "B", 16)
            pdf.multi_cell(w(), 8, ascii_pdf(title))
            pdf.ln(2)
            pdf.set_font("Helvetica", "", 11)
            for b in bullets:
                line = ascii_pdf("- " + b)
                pdf.multi_cell(w(), 5.5, line)
                pdf.ln(1)
        elif kind == "two_col":
            _, title, d = item
            pdf.set_font("Helvetica", "B", 16)
            pdf.multi_cell(w(), 8, ascii_pdf(title))
            pdf.ln(2)
            pdf.set_font("Helvetica", "B", 12)
            pdf.multi_cell(w(), 6, ascii_pdf(d["left_h"]))
            pdf.set_font("Helvetica", "", 10)
            for b in d["left"]:
                pdf.multi_cell(w(), 5, ascii_pdf("- " + b))
            pdf.ln(2)
            pdf.set_font("Helvetica", "B", 12)
            pdf.multi_cell(w(), 6, ascii_pdf(d["right_h"]))
            pdf.set_font("Helvetica", "", 10)
            for b in d["right"]:
                pdf.multi_cell(w(), 5, ascii_pdf("- " + b))

    pdf.output(str(OUT_PDF))


def main():
    slides = slide_plan()
    build_pptx(slides)
    build_pdf(slides)
    print(f"Wrote {OUT_PPTX}")
    print(f"Wrote {OUT_PDF} (open in Preview or view on GitHub.com)")


if __name__ == "__main__":
    main()
