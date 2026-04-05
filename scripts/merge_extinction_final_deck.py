#!/usr/bin/env python3
"""
Consolidate into Archive_of_Extinction_Final_BDC2026.pptx:
- Part I: BDC_Extinction_Archive_EN.pptx (theme, dimensions, in-place text fixes)
- Part II: Selected slides from BDC_Summit_Extinction_Archive_2026.pptx (summit / appendix text)
- Part III: Full-bleed images from BDC_Deck_EN.pptx

Extinction_Archive_Summit_2026.pptx is identical to the Summit file — not read twice.
"""

from __future__ import annotations

import re
import shutil
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
EXPORT = ROOT / "slides" / "export"
TEMPLATE = EXPORT / "BDC_Extinction_Archive_EN.pptx"
SUMMIT = EXPORT / "BDC_Summit_Extinction_Archive_2026.pptx"
IMAGE_DECK = EXPORT / "BDC_Deck_EN.pptx"
OUT = EXPORT / "Archive_of_Extinction_Final_BDC2026.pptx"

# 0-based summit slide indices to append (skip pure duplicates of Part I)
SUMMIT_APPEND_INDICES = [
    0,  # title stack
    1,  # Team
    2,  # Why this title
    6,  # Core design question
    7,  # Bio x digital quadrants
    8,  # Two dossiers
    9,  # Scene IDs
    10, # Design techniques
    11, # Ideation 8 steps
    12, # Four-week spine
    13, # Colonial / Country
    14, # Sherkow & Greely
    15, # Lit mammoth
    16, # Lit thylacine
    17, # Rubric
    19, # Repository pointers
]


def emu_to_inches(emu: int) -> float:
    return float(emu) / 914400.0


def replace_in_all_text_shapes(slide, mapping: list[tuple[str, str]]) -> None:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            full = "".join(run.text for run in para.runs)
            if not full.strip():
                continue
            new_full = full
            for old, new in mapping:
                if old in new_full:
                    new_full = new_full.replace(old, new)
            if new_full != full:
                for i, run in enumerate(para.runs):
                    run.text = new_full if i == 0 else ""


def slide_text_blob(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text and shape.text.strip():
            parts.append(shape.text.strip())
    return "\n\n".join(parts)


def add_blank_text_slide(prs, title: str, body_lines: list[str]) -> None:
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)
    margin = Inches(0.65)
    tw = prs.slide_width - 2 * margin
    th = prs.slide_height - 2 * margin
    box = slide.shapes.add_textbox(margin, margin, tw, th)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    for line in body_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.level = 0
        p.space_after = Pt(3)


def extract_slide_image_paths(pptx_path: Path) -> list[str]:
    with zipfile.ZipFile(pptx_path) as z:
        paths = []
        for i in range(1, 500):
            rels = f"ppt/slides/_rels/slide{i}.xml.rels"
            try:
                data = z.read(rels).decode("utf-8")
            except KeyError:
                break
            m = re.search(r'Target="([^"]+\.(png|jpg|jpeg))"', data, re.I)
            paths.append("" if not m else m.group(1).replace("../", "ppt/"))
        return paths


def merge() -> None:
    shutil.copy2(TEMPLATE, OUT)
    prs = Presentation(str(OUT))

    # Part I: coherent species + web-first physical + timeline
    replace_in_all_text_shapes(
        prs.slides[4],
        [
            ("Mammoth-first, pigeon as structural echo", "Mammoth + thylacine dossiers (two deep case studies)"),
            ("Passenger pigeon", "Thylacine (Tasmanian tiger)"),
            (
                "Mass flock synchrony\nDensity thresholds\nSocial collapse\nSky-darkening scale\nThinness as silence",
                "Orbit / diel activity (Pozniak 2018)\nPOV metaphor via RGC models (Mass & Supin 2020, Interpolated)\nColonial extinction history (Paddle; Sleightholme & Campbell)\nPalawa / Country framing (Rimmer; Lehman; Clements; Schlunke)\nEthics Console + reflection log",
            ),
        ],
    )

    replace_in_all_text_shapes(
        prs.slides[7],
        [
            ("Archive object + living bench", "Minimal physical anchor + WebXR entry (Summit build)"),
            (
                "3D-printed focal object anchors the body of the extinct species in gallery space.\nMobile WebAR opens dossier content from the physical object.\nSafe tactile materials symbolize tundra and steppe without claiming literal reconstruction.\nPlants plus soil moisture, light, and temperature sensors drive one live audio bus.",
                "3D-printed or surrogate focal object anchors attention in gallery space.\nQR / URL opens the WebXR memorial (A-Frame / Three.js).\nWeb-first scope: no living biosensor demo in v1; one citation one-sheet for judges.\nOptional WebAR from the printed anchor.",
            ),
            (
                "Not a simulation of extinct habitat — an index of present life.",
                "Not a simulation of extinct habitat — an index of evidence, uncertainty, and responsibility.",
            ),
        ],
    )

    replace_in_all_text_shapes(
        prs.slides[10],
        [
            (
                "Integrate passenger pigeon, ethics logic, WebAR, tactile table, and living sensor pipeline.",
                "Integrate thylacine dossier, Indigenous context UI, ethics fork (Sherkow & Greely), reflection log, WebXR polish.",
            ),
        ],
    )

    replace_in_all_text_shapes(
        prs.slides[0],
        [
            (
                "Literature-grounded memorial · biodigital experience",
                "Literature-grounded memorial · biodigital experience\nMacau University · Mentor: Atticus SIMS · GUO XIAO YUE (MC569254) · LIU JIA QUN (MC569293)",
            ),
        ],
    )

    add_blank_text_slide(
        prs,
        "Part II — Summit appendix (English)",
        [
            "Selected slides from BDC_Summit_Extinction_Archive_2026: citations, scene IDs, rubric, methods.",
            "Part I (previous slides) follows the BDC_Extinction_Archive_EN layout and narrative spine.",
        ],
    )

    summit_prs = Presentation(str(SUMMIT))
    for idx in SUMMIT_APPEND_INDICES:
        if idx >= len(summit_prs.slides):
            continue
        blob = slide_text_blob(summit_prs.slides[idx])
        if not blob.strip():
            continue
        lines = [ln.strip() for ln in blob.split("\n") if ln.strip()]
        title_line = lines[0][:200] if lines else f"Slide {idx+1}"
        body = lines[1:22]
        if not body:
            body = [blob[:800]]
        add_blank_text_slide(prs, title_line, body)

    if IMAGE_DECK.exists():
        add_blank_text_slide(
            prs,
            "Part III — Visual deck (BDC_Deck_EN)",
            [
                "Full-bleed slide graphics from BDC_Deck_EN.pptx (image export pipeline).",
            ],
        )
        media_paths = extract_slide_image_paths(IMAGE_DECK)
        sw = emu_to_inches(prs.slide_width)
        sh = emu_to_inches(prs.slide_height)
        with tempfile.TemporaryDirectory() as td:
            tdp = Path(td)
            with zipfile.ZipFile(IMAGE_DECK) as z:
                z.extractall(td)
            for rel in media_paths:
                if not rel:
                    continue
                abs_path = tdp / rel
                if not abs_path.is_file():
                    continue
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                try:
                    slide.shapes.add_picture(str(abs_path), Inches(0), Inches(0), width=Inches(sw), height=Inches(sh))
                except OSError:
                    pass

    add_blank_text_slide(
        prs,
        "Thank you",
        [
            "Archive of Extinction — Umwelt Hypothesis Dossiers",
            "Macau University · Biodesign Challenge 2026 · Biodigital Excellence",
            "GUO XIAO YUE (MC569254) · LIU JIA QUN (MC569293) · Mentor: Atticus SIMS",
        ],
    )

    prs.save(str(OUT))
    print(f"Wrote {OUT} — {len(prs.slides)} slides")


if __name__ == "__main__":
    merge()
